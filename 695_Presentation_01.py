from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Deep-Sea Fish Detection Using CNNs"
subtitle.text = "A Machine Learning Approach\n[Your Name] | [Date] | [Institution]"

# Slide 2: Project Overview
slide_layout = prs.slide_layouts[1]  # Title and content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Project Overview"
content.text = ("• Detecting deep-sea fish species is challenging due to low visibility and complex backgrounds.\n"
                "• CNNs provide a powerful tool for accurate species detection and classification.\n"
                "• Goal: Develop a robust CNN model for real-world underwater detection.")

# Slide 3: Methodology
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Methodology: Convolutional Neural Networks (CNNs)"
content.text = ("• Base Models: VGG16, ResNet50, EfficientNet\n"
                "• Data Augmentation: Rotation, scaling, contrast adjustments\n"
                "• Transfer Learning: Using pre-trained models for improved accuracy\n"
                "• Hyperparameter Optimization: Fine-tuning learning rate, batch size, dropout")

# Slide 4: Dataset - SeamapD21
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Dataset: SEAMAPD21"
content.text = ("• Publicly available dataset with ~28,328 images & 90,000 annotations\n"
                "• 130 fish species captured in real-world underwater environments\n"
                "• Includes complex backgrounds and varying lighting conditions\n"
                "• Source: GitHub (https://github.com/SEFSC/SEAMAPD21)")

# Slide 5: Initial Findings
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Initial Findings"
content.text = ("• Literature review highlights success with YOLO and CNN-based models\n"
                "• Key challenge: Limited annotated datasets for deep-sea species\n"
                "• Need for robust preprocessing to enhance model performance")

# Slide 6: Challenges & Solutions
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Challenges and Proposed Solutions"
content.text = ("• Dataset Limitations: Need more annotated images → Solution: Collaborative data collection\n"
                "• Low Visibility & Background Noise → Solution: Advanced image enhancement\n"
                "• Similar-Looking Species → Solution: Feature engineering and deeper CNN layers")

# Slide 7: Next Steps (2-Week Plan)
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Next Steps: Two-Week Plan"
content.text = ("Week 1:\n"
                "• Data exploration and preprocessing\n"
                "• Implement augmentation and prepare dataset\n"
                "• Train baseline CNN models\n\n"
                "Week 2:\n"
                "• Optimize CNN architectures (VGG16, ResNet, EfficientNet)\n"
                "• Evaluate performance (precision, recall, F1-score)\n"
                "• Analyze misclassifications and refine approach")

# Slide 8: Conclusion
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusion"
content.text = ("• CNNs offer a promising solution for deep-sea fish detection\n"
                "• Challenges include dataset limitations and underwater image complexity\n"
                "• Next steps focus on refining models for better accuracy\n"
                "• Expected impact: Improved marine species monitoring and conservation efforts")

# Save the presentation file
#ppt_filename = "/mnt/data/Deep_Sea_Fish_Detection_Presentation.pptx"
ppt_filename = r"C:\Users\w10185657\OneDrive - The University of Southern Mississippi\695_Project_Thesis\Deep_Sea_Fish_Detection_Presentation.pptx"
prs.save(ppt_filename)

ppt_filename
